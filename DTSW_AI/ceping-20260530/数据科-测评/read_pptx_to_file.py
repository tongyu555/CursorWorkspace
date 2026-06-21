import sys
from pptx import Presentation

def extract_text(pptx_path, out_path):
    prs = Presentation(pptx_path)
    with open(out_path, "w", encoding="utf-8") as f:
        for i, slide in enumerate(prs.slides):
            f.write(f"--- Slide {i+1} ---\n")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    f.write(shape.text + "\n")
            f.write("\n")

if __name__ == "__main__":
    extract_text("e:/CursorWorkspace/DTSW_AI/ceping-20260530/数据科-测评/站址合并算法说明.pptx", "e:/CursorWorkspace/DTSW_AI/ceping-20260530/数据科-测评/pptx_text.txt")
