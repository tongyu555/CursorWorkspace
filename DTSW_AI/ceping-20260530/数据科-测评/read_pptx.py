import sys
from pptx import Presentation

def extract_text(pptx_path):
    prs = Presentation(pptx_path)
    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)
        print()

if __name__ == "__main__":
    extract_text("e:/CursorWorkspace/DTSW_AI/ceping-20260530/数据科-测评/站址合并算法说明.pptx")
